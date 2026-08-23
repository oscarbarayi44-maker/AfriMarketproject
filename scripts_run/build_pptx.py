"""Génère docs/AfriMarket_Presentation_Direction.pptx — présentation
professionnelle 16:9 suivant la charte graphique AfriMarket."""
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

PROJECT_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(PROJECT_ROOT / "src"))
from afrimarket import branding  # noqa: E402

FIGS = PROJECT_ROOT / "figures_brand"
OUT = PROJECT_ROOT / "docs" / "AfriMarket_Presentation_Direction.pptx"
RESULTS = json.loads((PROJECT_ROOT / "docs" / "resultats_analyses.json").read_text(encoding="utf-8"))

DARK_GREEN = RGBColor.from_string(branding.DARK_GREEN.lstrip("#"))
BRAND_GREEN = RGBColor.from_string(branding.BRAND_GREEN.lstrip("#"))
LIGHT_GREEN = RGBColor.from_string(branding.LIGHT_GREEN.lstrip("#"))
PALE_GREEN = RGBColor.from_string(branding.PALE_GREEN.lstrip("#"))
WHITE = RGBColor.from_string("FFFFFF")
TEXT_DARK = RGBColor.from_string(branding.TEXT_DARK.lstrip("#"))
GREY = RGBColor.from_string(branding.GREY.lstrip("#"))

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line:
        shape.line.color.rgb = color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def textbox(slide, x, y, w, h, text, size=18, color=TEXT_DARK, bold=False, italic=False,
            align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        run.font.color.rgb = color
    return tb


def bullets(slide, x, y, w, h, items, size=15, color=TEXT_DARK, bullet_color=None, font="Calibri", space_after=10):
    bullet_color = bullet_color or BRAND_GREEN
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = f"●  {item}"
        run.font.size = Pt(size)
        run.font.name = font
        run.font.color.rgb = color
    return tb


def footer(slide, page_no):
    rect(slide, 0, SLIDE_H - Inches(0.32), SLIDE_W, Inches(0.32), DARK_GREEN)
    textbox(slide, Inches(0.4), SLIDE_H - Inches(0.34), Inches(6), Inches(0.32),
            "AfriMarket — Analyse stratégique", size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    textbox(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.34), Inches(0.8), Inches(0.32),
            str(page_no), size=10, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, SLIDE_W, Inches(1.15), DARK_GREEN)
    rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.06), BRAND_GREEN)
    if branding.LOGO_PATH:
        slide.shapes.add_picture(str(branding.LOGO_PATH), SLIDE_W - Inches(1.15), Inches(0.12), height=Inches(0.9))
    textbox(slide, Inches(0.5), Inches(0.18), Inches(10.5), Inches(0.6), title, size=26, color=WHITE, bold=True)
    if subtitle:
        textbox(slide, Inches(0.5), Inches(0.7), Inches(10.5), Inches(0.4), subtitle, size=13, color=LIGHT_GREEN)


def picture_fit(slide, path, x, y, max_w, max_h):
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    ratio = min(max_w / iw, max_h / ih)
    w, h = Emu(int(iw * ratio)), Emu(int(ih * ratio))
    px = x + (max_w - w) // 2
    py = y + (max_h - h) // 2
    slide.shapes.add_picture(str(path), px, py, width=w, height=h)


# ============================================================ SLIDE 1 : COVER
s = add_slide()
rect(s, 0, 0, SLIDE_W, SLIDE_H, DARK_GREEN)
rect(s, 0, Inches(4.55), SLIDE_W, Inches(0.08), BRAND_GREEN)
if branding.LOGO_PATH:
    s.shapes.add_picture(str(branding.LOGO_PATH), Inches(5.67), Inches(1.1), height=Inches(1.8))
    title_y = Inches(3.15)
else:
    textbox(s, Inches(0), Inches(1.5), SLIDE_W, Inches(1.2), "AfriMarket", size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    title_y = Inches(2.9)
textbox(s, Inches(0.8), title_y, SLIDE_W - Inches(1.6), Inches(1.0),
        "Analyse stratégique des données commerciales", size=30, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
textbox(s, Inches(0.8), title_y + Inches(0.85), SLIDE_W - Inches(1.6), Inches(0.5),
        "Résultats, enseignements clés et recommandations — 6 mois d'activité",
        size=16, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)
textbox(s, Inches(0.8), SLIDE_H - Inches(1.1), SLIDE_W - Inches(1.6), Inches(0.5),
        "Présenté à la Direction Générale  |  Data Analyst", size=13, color=WHITE, align=PP_ALIGN.CENTER)

# ============================================================ SLIDE 2 : SOMMAIRE
s = add_slide()
header(s, "Sommaire")
items = [
    "1.  Contexte & méthodologie",
    "2.  Performance globale",
    "3.  Analyse par catégorie — quelle catégorie prioriser ?",
    "4.  Analyse géographique — où investir ?",
    "5.  Analyse marketing — quel canal renforcer ?",
    "6.  Analyse clients — comment fidéliser ?",
    "7.  5 recommandations stratégiques",
    "8.  Conclusion & plan d'action",
]
bullets(s, Inches(1.2), Inches(1.7), Inches(10.9), Inches(5), items, size=20, space_after=16)
footer(s, 2)

# ============================================================ SLIDE 3 : CONTEXTE & METHODO
s = add_slide()
header(s, "Contexte & méthodologie", "6 mois d'activité commerciale — 4 catégories, 8 villes")
textbox(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.8),
        "AfriMarket constate des variations de CA, un taux de retour préoccupant, des dépenses marketing "
        "élevées et des écarts de performance selon les villes. Cette analyse répond à 4 questions "
        "stratégiques business à partir d'un audit rigoureux des données.", size=14, color=TEXT_DARK)

rect(s, Inches(0.5), Inches(2.5), Inches(5.9), Inches(4.3), PALE_GREEN)
textbox(s, Inches(0.75), Inches(2.7), Inches(5.4), Inches(0.4), "Qualité des données brutes", size=16, bold=True, color=DARK_GREEN)
bullets(s, Inches(0.75), Inches(3.25), Inches(5.4), Inches(3.4), [
    "10 100 commandes brutes analysées",
    "100 doublons détectés et supprimés",
    "614 remises négatives corrigées",
    "632 prix négatifs / 799 valeurs extrêmes traités",
    "608 quantités nulles supprimées",
    "Villes, catégories et statuts uniformisés",
], size=13)

rect(s, Inches(6.9), Inches(2.5), Inches(5.9), Inches(4.3), PALE_GREEN)
textbox(s, Inches(7.15), Inches(2.7), Inches(5.4), Inches(0.4), "Dataset final exploitable", size=16, bold=True, color=DARK_GREEN)
bullets(s, Inches(7.15), Inches(3.25), Inches(5.4), Inches(3.4), [
    "9 400 commandes valides (df_clean)",
    "1 734 clients uniques",
    "Hypothèse retenue : marge brute = 35% du CA (à défaut de coût matière communiqué)",
    "Indicateurs recalculés : CA, marge, profit net, taux de retour, CLV",
], size=13)
footer(s, 3)

# ============================================================ SLIDE 4 : PERFORMANCE GLOBALE
perf = RESULTS["performance_globale"]
s = add_slide()
header(s, "Performance globale", "Vue d'ensemble sur 6 mois")
kpis = [
    ("CA total", f"{perf['ca_total']/1e6:.2f} M$"),
    ("Profit net estimé", f"{perf['profit_net_total']/1e3:.0f} k$"),
    ("Panier moyen", f"{perf['panier_moyen']:.0f} $"),
    ("Taux d'annulation", f"{perf['taux_annulation']*100:.1f} %"),
    ("Taux de retour", f"{perf['taux_retour']*100:.1f} %"),
]
card_w = Inches(2.35)
gap = Inches(0.25)
start_x = Inches(0.5)
for i, (label, value) in enumerate(kpis):
    x = start_x + i * (card_w + gap)
    rect(s, x, Inches(2.2), card_w, Inches(2.4), DARK_GREEN if i == 0 else BRAND_GREEN)
    textbox(s, x, Inches(2.45), card_w, Inches(1.1), value, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(s, x + Inches(0.1), Inches(3.6), card_w - Inches(0.2), Inches(0.9), label, size=13, color=WHITE, align=PP_ALIGN.CENTER)
textbox(s, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.6),
        "Sur 6 mois, AfriMarket génère 2,51 M$ de chiffre d'affaires pour un profit net estimé de 754 k$ "
        "(≈30% de marge nette). Le taux d'annulation reste faible (1,9%), mais le taux de retour (8,3%) "
        "constitue un point de vigilance direct sur la rentabilité.", size=14, color=TEXT_DARK)
footer(s, 4)


def analysis_slide(page_no, title, subtitle, image_name, insight_title, insight_text, image_w=Inches(7.6)):
    s = add_slide()
    header(s, title, subtitle)
    picture_fit(s, FIGS / image_name, Inches(0.4), Inches(1.4), image_w, Inches(5.6))
    tx = Inches(0.4) + image_w + Inches(0.3)
    tw = SLIDE_W - tx - Inches(0.4)
    rect(s, tx, Inches(1.4), tw, Inches(5.6), PALE_GREEN)
    textbox(s, tx + Inches(0.25), Inches(1.6), tw - Inches(0.5), Inches(0.5), insight_title, size=15, bold=True, color=DARK_GREEN)
    textbox(s, tx + Inches(0.25), Inches(2.15), tw - Inches(0.5), Inches(4.6), insight_text, size=13, color=TEXT_DARK, line_spacing=1.15)
    footer(s, page_no)
    return s


# ============================================================ SLIDE 5 : CATEGORIE
analysis_slide(
    5, "Analyse par catégorie", "Question stratégique : quelle catégorie prioriser ou optimiser ?",
    "ca_par_categorie.png",
    "Constat",
    "L'Électronique porte 74,6% du CA (1,87 M$) et 82% du profit net.\n\n"
    "Mode et Beauté ont un profit quasi nul : les coûts logistiques et "
    "marketing absorbent 57% et 88% de leur marge brute (contre 6% pour "
    "l'Électronique).\n\n"
    "Recommandation : prioriser l'Électronique, optimiser Mode/Beauté par "
    "le panier moyen plutôt que par le volume.",
)

# ============================================================ SLIDE 6 : CATEGORIE RETOUR
analysis_slide(
    6, "Analyse par catégorie", "Le revers de la médaille : le taux de retour",
    "taux_retour_categorie.png",
    "Point de vigilance",
    "Le taux de retour de l'Électronique atteint 14,1%, près de 2× la "
    "moyenne globale (8,3%).\n\n"
    "Chaque point de retour gagné se traduit directement en profit "
    "préservé sur la catégorie la plus rentable.\n\n"
    "Actions : fiches produits plus précises, contrôle qualité "
    "fournisseurs, politique de retour resserrée sur les motifs évitables.",
)

# ============================================================ SLIDE 7 : GEOGRAPHIE
analysis_slide(
    7, "Analyse géographique", "Question stratégique : où devons-nous investir davantage ?",
    "ca_par_ville.png",
    "Constat",
    "Kinshasa est le marché le plus performant (752 k$ de CA, 227 k$ de "
    "profit) avec un taux d'annulation quasi nul (0,26%).\n\n"
    "Abidjan et Dakar combinent bon CA et annulation quasi nulle : "
    "marchés secondaires à consolider.\n\n"
    "Recommandation : renforcer l'investissement à Kinshasa, consolider "
    "Abidjan et Dakar.",
)

# ============================================================ SLIDE 8 : GEOGRAPHIE ANOMALIE
analysis_slide(
    8, "Analyse géographique", "Une anomalie opérationnelle à traiter en priorité",
    "taux_annulation_ville.png",
    "Alerte",
    "Douala affiche un taux d'annulation de 12,9%, contre moins de 1% "
    "dans toutes les autres villes — une anomalie opérationnelle majeure.\n\n"
    "Causes probables à auditer : méthode de paiement locale, rupture de "
    "stock, fiabilité du livreur partenaire.\n\n"
    "Recommandation : auditer la cause racine avant tout investissement "
    "marketing supplémentaire sur cette ville.",
)

# ============================================================ SLIDE 9 : MARKETING
analysis_slide(
    9, "Analyse marketing", "Question stratégique : quel canal mérite plus de budget ?",
    "roi_par_canal.png",
    "ROI = (Revenus − Coût marketing) / Coût marketing",
    "Email affiche le meilleur ROI (230) pour le plus petit budget "
    "(2 297 $) : canal nettement sous-financé.\n\n"
    "Google Ads (ROI 50) reste solide et scalable.\n\n"
    "Influenceur cumule le pire ROI (21,5) pour un budget de 16 k$.\n\n"
    "Recommandation : augmenter Email et Google Ads, réduire Influenceur.",
)

# ============================================================ SLIDE 10 : MARKETING RETENTION
analysis_slide(
    10, "Analyse marketing", "La rétention client par canal, en complément du ROI",
    "retention_par_canal.png",
    "Constat",
    "Instagram Ads génère le plus de CA (950 k$) et la meilleure "
    "rétention (54%), malgré un ROI plus modéré (25) lié à son budget "
    "élevé (37 k$).\n\n"
    "Influenceur a aussi la pire rétention (42%) : canal à réduire ou "
    "réattribuer en priorité.\n\n"
    "Recommandation : conserver Instagram Ads pour le volume et la "
    "fidélisation, en optimisant le ciblage plutôt qu'en augmentant le "
    "budget brut.",
)

# ============================================================ SLIDE 11 : CLIENTS
analysis_slide(
    11, "Analyse clients", "Question stratégique : comment améliorer la rétention ?",
    "pareto_clients.png",
    "Constat",
    "1 734 clients, dont 73,2% sont récurrents (≥2 commandes) : un "
    "socle de fidélité globalement sain.\n\n"
    "31,9% des clients génèrent 80% du CA — une concentration modérée, "
    "plutôt rassurante par rapport à un Pareto classique.\n\n"
    "Recommandation : capitaliser sur cette fidélité existante par une "
    "segmentation différenciée (slide suivante).",
)

# ============================================================ SLIDE 12 : SEGMENTATION
analysis_slide(
    12, "Analyse clients", "Segmentation par valeur vie client (quartiles)",
    "segmentation_clients.png",
    "Constat",
    "Le segment Platine (25% des clients) génère 72% du CA total "
    "(1,81 M$), contre seulement 1,4% pour le segment Bronze.\n\n"
    "Une forte polarisation de la valeur client, à exploiter par une "
    "stratégie différenciée plutôt qu'uniforme.\n\n"
    "Recommandation : programme VIP pour Platine, offres d'activation "
    "pour Bronze/Argent (867 clients, 7,8% du CA).",
    image_w=Inches(5.6),
)

# ============================================================ SLIDE 13 : RECOMMANDATIONS
s = add_slide()
header(s, "5 recommandations stratégiques")
recos = [
    ("1", "Prioriser l'Électronique, réduire son taux de retour",
     "74,6% du CA et 82% du profit ; agir sur les fiches produits et le contrôle qualité fournisseurs."),
    ("2", "Revoir l'économie unitaire de Mode et Beauté",
     "Coûts logistiques/marketing absorbant 57%/88% de la marge ; seuil de livraison gratuite, bundles."),
    ("3", "Auditer en urgence l'anomalie de Douala",
     "12,9% d'annulation vs <1% ailleurs ; identifier la cause racine avant d'investir davantage."),
    ("4", "Réallouer le budget marketing vers l'efficacité",
     "Renforcer Email (ROI 230) et Google Ads (ROI 50) ; réduire Influenceur (ROI 21,5)."),
    ("5", "Structurer un programme de fidélisation par segment",
     "VIP pour Platine (72% du CA) ; offres d'activation pour Bronze/Argent (7,8% du CA)."),
]
y = Inches(1.45)
for num, title, desc in recos:
    rect(s, Inches(0.5), y, Inches(0.7), Inches(0.7), DARK_GREEN)
    textbox(s, Inches(0.5), y, Inches(0.7), Inches(0.7), num, size=24, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(1.4), y - Inches(0.02), Inches(11.2), Inches(0.4), title, size=16, bold=True, color=DARK_GREEN)
    textbox(s, Inches(1.4), y + Inches(0.38), Inches(11.2), Inches(0.6), desc, size=12.5, color=TEXT_DARK)
    y += Inches(1.08)
footer(s, 13)

# ============================================================ SLIDE 14 : CONCLUSION
s = add_slide()
header(s, "Conclusion & plan d'action")
textbox(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.6),
        "AfriMarket dispose d'un socle commercial sain : 2,51 M$ de CA, 754 k$ de profit net estimé, une "
        "base client majoritairement fidèle (73,2% de récurrence) et un taux d'annulation faible (1,9%).",
        size=15, color=TEXT_DARK)
textbox(s, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.4), "3 priorités d'action immédiates :", size=16, bold=True, color=DARK_GREEN)
bullets(s, Inches(0.7), Inches(3.5), Inches(11.9), Inches(2.2), [
    "Consolider l'Électronique en maîtrisant son taux de retour",
    "Corriger la rentabilité structurelle de Mode et Beauté par le panier moyen",
    "Traiter sans délai l'anomalie opérationnelle de Douala",
], size=14, space_after=12)
textbox(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.2),
        "Un réarbitrage du budget marketing (moins d'Influenceur, plus d'Email et Google Ads) peut "
        "améliorer le ROI global sans augmenter la dépense totale. Ces indicateurs sont pilotables "
        "mensuellement dans le dashboard Streamlit livré avec cette analyse.",
        size=13, italic=True, color=GREY)
footer(s, 14)

# ============================================================ SLIDE 15 : MERCI
s = add_slide()
rect(s, 0, 0, SLIDE_W, SLIDE_H, DARK_GREEN)
if branding.LOGO_PATH:
    s.shapes.add_picture(str(branding.LOGO_PATH), Inches(5.92), Inches(1.6), height=Inches(1.4))
textbox(s, Inches(0), Inches(3.3), SLIDE_W, Inches(1.0), "Merci", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(s, Inches(0), Inches(4.2), SLIDE_W, Inches(0.6), "Questions & discussion", size=18, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)

OUT.parent.mkdir(exist_ok=True)
prs.save(OUT)
print(f"Présentation écrite : {OUT}")
